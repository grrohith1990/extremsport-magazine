#!/usr/bin/env python3
"""
Enhanced Case Study PowerPoint: Extremsport Online-Magazin
Professional presentation with speaker notes for interview
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# === Color Palette ===
BG_DARK = RGBColor(0x0F, 0x0F, 0x1A)
ACCENT_RED = RGBColor(0xE7, 0x4C, 0x3C)
ACCENT_ORANGE = RGBColor(0xF3, 0x9C, 0x12)
ACCENT_TEAL = RGBColor(0x1A, 0xBC, 0x9C)
ACCENT_BLUE = RGBColor(0x34, 0x98, 0xDB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xEC, 0xF0, 0xF1)
SUBTLE = RGBColor(0x99, 0x99, 0xAA)
MID = RGBColor(0x66, 0x66, 0x77)
DARK_TEXT = RGBColor(0x2C, 0x3E, 0x50)
GREEN = RGBColor(0x27, 0xAE, 0x60)
PURPLE = RGBColor(0x9B, 0x59, 0xB6)


def set_bg(slide, color=BG_DARK):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def txt(slide, left, top, width, height, text, size=18, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, font="Segoe UI", italic=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font
    p.font.italic = italic
    p.alignment = align
    return box


def multi_txt(slide, left, top, width, height, lines, size=16, color=WHITE, spacing=Pt(8)):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, (text, fmt) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(fmt.get("size", size))
        p.font.color.rgb = fmt.get("color", color)
        p.font.bold = fmt.get("bold", False)
        p.font.italic = fmt.get("italic", False)
        p.font.name = "Segoe UI"
        p.space_after = fmt.get("spacing", spacing)
        p.alignment = fmt.get("align", PP_ALIGN.LEFT)
    return box


def accent_bar(slide, left=Inches(0.8), top=Inches(1.35), width=Inches(2.5)):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.04))
    s.fill.solid()
    s.fill.fore_color.rgb = ACCENT_RED
    s.line.fill.background()


def side_accent(slide):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.08), Inches(7.5))
    s.fill.solid()
    s.fill.fore_color.rgb = ACCENT_RED
    s.line.fill.background()


def box(slide, left, top, width, height, text, fill, text_color=WHITE, size=11, sub=""):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = fill
    s.shadow.inherit = False
    tf = s.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = text_color
    p.font.bold = True
    p.font.name = "Segoe UI"
    p.alignment = PP_ALIGN.CENTER
    if sub:
        p2 = tf.add_paragraph()
        p2.text = sub
        p2.font.size = Pt(size - 2)
        p2.font.color.rgb = text_color
        p2.font.name = "Segoe UI"
        p2.alignment = PP_ALIGN.CENTER
    return s


def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def slide_header(slide, title, subtitle=""):
    set_bg(slide, BG_DARK)
    side_accent(slide)
    txt(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.9), title, size=34, bold=True, color=WHITE)
    accent_bar(slide)
    if subtitle:
        txt(slide, Inches(0.8), Inches(1.55), Inches(11), Inches(0.6), subtitle, size=14, color=SUBTLE, italic=True)


# ═══════════════════════════════════════════════════════════
# SLIDE 1: TITLE
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)

s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.12), Inches(7.5))
s.fill.solid()
s.fill.fore_color.rgb = ACCENT_RED
s.line.fill.background()

for cx, cy, sz in [(11.5, 1.0, 2.5), (12.0, 5.5, 1.8), (10.0, 6.5, 1.2)]:
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx), Inches(cy), Inches(sz), Inches(sz))
    c.fill.solid()
    c.fill.fore_color.rgb = ACCENT_RED
    c.fill.fore_color.brightness = 0.7
    c.line.fill.background()

txt(slide, Inches(1.0), Inches(1.2), Inches(4), Inches(0.5),
    "CASE STUDY  ·  ARCHITEKTUR", size=12, bold=True, color=ACCENT_RED)
txt(slide, Inches(1.0), Inches(2.0), Inches(10), Inches(2.0),
    "Extremsport\nOnline-Magazin", size=52, bold=True, color=WHITE, font="Segoe UI Light")
txt(slide, Inches(1.0), Inches(4.2), Inches(9), Inches(1.0),
    "Agile Architektur für ein modernes, skalierbares Medienprodukt\nmit austauschbaren Komponenten und Cloud-nativer Infrastruktur",
    size=17, color=SUBTLE)
txt(slide, Inches(1.0), Inches(6.5), Inches(8), Inches(0.5),
    "Juli 2026  ·  20 Minuten  ·  Freie Medienwahl", size=13, color=MID)

add_notes(slide, """EINSTIEG (1 Min):
- Kurz vorstellen, dann direkt ins Thema
- "Ich zeige heute, wie ich ein Extremsport-Onlinemagazin architektonisch aufsetzen wuerde - mit besonderem Fokus auf Agilitaet und Veraenderbarkeit."
- Energie zeigen! Extremsport = Leidenschaft
- Blickkontakt, selbstbewusst starten""")


# ═══════════════════════════════════════════════════════════
# SLIDE 2: AGENDA
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Agenda")

items = [
    ("01", "Ausgangslage & Herausforderungen", "Was der Kunde braucht"),
    ("02", "Qualitaetsanforderungen (NFRs)", "Messbare Architekturziele"),
    ("03", "Architekturentscheidungen", "Warum Microservices + Adapter"),
    ("04", "High-Level Architektur", "System mit Umsystemen"),
    ("05", "Evolutionspfad & Risiken", "Wie die Architektur agil mitwaechst"),
]

y = Inches(2.2)
for num, title, desc in items:
    txt(slide, Inches(1.0), y, Inches(0.8), Inches(0.5), num, size=28, bold=True, color=ACCENT_RED, font="Segoe UI Light")
    txt(slide, Inches(2.0), y, Inches(6), Inches(0.4), title, size=19, bold=True, color=WHITE)
    txt(slide, Inches(2.0), y + Inches(0.38), Inches(6), Inches(0.4), desc, size=13, color=SUBTLE)
    y += Inches(0.95)

add_notes(slide, """UEBERLEITUNG (30 Sek):
- "Ich starte mit der Ausgangslage, leite daraus Qualitaetsanforderungen ab, zeige dann meine Architekturentscheidungen und wie das Ganze zusammenspielt."
- Zeigt strukturiertes Denken""")


# ═══════════════════════════════════════════════════════════
# SLIDE 3: AUSGANGSLAGE
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Ausgangslage", "Was gebaut werden soll + existierende Systemlandschaft")

txt(slide, Inches(0.8), Inches(2.0), Inches(5), Inches(0.4), "NEUES SYSTEM", size=11, bold=True, color=ACCENT_TEAL)

new_items = [
    ("Kundenportal", "Free + Premium (Abo-Modell)"),
    ("Artikel-Bereich", "Einzel-Kauf & Abo-Zugang"),
    ("Moderiertes Forum", "Community-Building"),
    ("Autorenbereich", "Content-Erstellung & Verwaltung"),
    ("Responsive Frontend", "Mobile First / PWA"),
    ("Archivsystem", "Aeltere Inhalte durchsuchbar"),
]
y = Inches(2.4)
for title, desc in new_items:
    txt(slide, Inches(1.0), y, Inches(3.5), Inches(0.35), f"▸  {title}", size=15, bold=True, color=WHITE)
    txt(slide, Inches(1.5), y + Inches(0.3), Inches(4), Inches(0.3), desc, size=12, color=SUBTLE)
    y += Inches(0.65)

txt(slide, Inches(7.0), Inches(2.0), Inches(5), Inches(0.4), "UMSYSTEME (BESTAND)", size=11, bold=True, color=ACCENT_ORANGE)

ext_items = [
    ("CMS", "IBM FileNet -> Cloud-CMS?", ACCENT_ORANGE),
    ("Buchhaltung", "Win2000 - Legacy, bleibt", ACCENT_RED),
    ("Merch-Shop", "Drittanbieter + Versand", ACCENT_BLUE),
    ("Auth-Server", "AD-Eigenbau -> Keycloak?", ACCENT_ORANGE),
]
y = Inches(2.5)
for title, desc, col in ext_items:
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.0), y + Inches(0.08), Inches(0.15), Inches(0.15))
    s.fill.solid()
    s.fill.fore_color.rgb = col
    s.line.fill.background()
    txt(slide, Inches(7.4), y, Inches(3), Inches(0.35), title, size=15, bold=True, color=WHITE)
    txt(slide, Inches(7.4), y + Inches(0.3), Inches(4.5), Inches(0.3), desc, size=12, color=SUBTLE)
    y += Inches(0.75)

box(slide, Inches(7.0), Inches(5.8), Inches(5.5), Inches(1.0),
    "⚡ Zwei Abloesungen unklar ->\nArchitektur MUSS Austauschbarkeit garantieren",
    RGBColor(0x2C, 0x1A, 0x1A), ACCENT_ORANGE, size=12)

add_notes(slide, """AUSGANGSLAGE (3 Min):
- Links: "Das sind die 6 Kernbausteine die wir neu bauen"
- Rechts: "Existierende Landschaft - und hier wird es spannend:"
- BETONEN: CMS und Auth sind UNSICHER ob sie abgeloest werden
- "Das ist DER zentrale Treiber fuer meine Architekturentscheidungen"
- Orange = Unsicherheit/Veraenderung geplant
- Rot = Legacy, kein Austausch
- Blau = Stabil, externe API
- Ueberleitung: "Daraus leite ich jetzt meine Qualitaetsanforderungen ab" """)


# ═══════════════════════════════════════════════════════════
# SLIDE 4: NFRs
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Qualitaetsanforderungen", "5 messbare NFRs als Architektur-Treiber")

nfrs = [
    ("Modifizierbarkeit", "Austausch eines Umsystems (CMS/Auth) in < 2 Sprints\nohne Aenderung an anderen Services", ACCENT_ORANGE, "★★★"),
    ("Skalierbarkeit", "10x Lastspitze (virale Events) innerhalb von 2 Min.\nautomatisch abfangen (Horizontal Scaling)", ACCENT_BLUE, "★★★"),
    ("Verfuegbarkeit", "99,5% Uptime Public-Bereich; Graceful Degradation\nbei Ausfall einzelner Services", ACCENT_TEAL, "★★☆"),
    ("Sicherheit", "DSGVO-konform; OAuth2/OIDC; Rechtetrennung\n(User / Autor / Moderator / Admin); Pen-Test ready", ACCENT_RED, "★★☆"),
    ("Performance", "Seitenladezeit < 2s (P95); API < 300ms;\nCore Web Vitals im gruenen Bereich", PURPLE, "★★☆"),
]

y = Inches(2.1)
for i, (title, desc, color, prio) in enumerate(nfrs):
    txt(slide, Inches(0.8), y, Inches(1), Inches(0.35), prio, size=14, color=color)
    txt(slide, Inches(1.8), y, Inches(3), Inches(0.35), title, size=17, bold=True, color=color)
    txt(slide, Inches(4.8), y, Inches(8), Inches(0.7), desc, size=13, color=LIGHT)
    y += Inches(0.95)
    if i < 4:
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), y - Inches(0.12), Inches(11.5), Inches(0.01))
        s.fill.solid()
        s.fill.fore_color.rgb = RGBColor(0x2A, 0x2A, 0x3E)
        s.line.fill.background()

txt(slide, Inches(0.8), Inches(6.8), Inches(11), Inches(0.4),
    "★★★ = Architektur-bestimmend  |  ★★☆ = Wichtig, beeinflusst Design-Details", size=11, color=MID)

add_notes(slide, """NFRs (3 Min):
- "5 Qualitaetsanforderungen, priorisiert nach Architektur-Impact"
- Modifizierbarkeit ZUERST: "Wenn wir in 6 Monaten das CMS tauschen, darf das nicht den Rest betreffen"
- Skalierbarkeit: "Extremsport ist eventgetrieben - viraler Klippensprung = 10x Traffic"
- Messbarkeit betonen: "Alle NFRs sind testbar formuliert"
- TIPP: Fragen ob jemand ergaenzen wuerde -> zeigt Teamfaehigkeit""")


# ═══════════════════════════════════════════════════════════
# SLIDE 5: ARCHITEKTURENTSCHEIDUNGEN
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Architekturentscheidungen", "Entscheidungslogik: Welches Pattern loest welches Problem?")

decisions = [
    ("Microservices", "Modifizierbarkeit\nSkalierbarkeit",
     "Unabhaengige Deployments;\nServices einzeln skalierbar;\nTeam-Autonomie",
     "Operational Complexity;\nEventual Consistency"),
    ("Adapter Pattern", "Modifizierbarkeit",
     "CMS/Auth austauschbar ohne\nKernlogik zu aendern;\nAnti-Corruption Layer",
     "Extra Abstraktionsschicht;\ninitialer Mehraufwand"),
    ("API Gateway", "Sicherheit\nPerformance",
     "Zentraler Entry Point;\nRate Limiting; Auth-Validierung;\nRouting & Load Balancing",
     "Single Point of Failure\n(-> Redundanz noetig)"),
    ("Event-Driven\n(async)", "Skalierbarkeit\nVerfuegbarkeit",
     "Asynchrone Verarbeitung;\nLoose Coupling;\nResilience bei Lastspitzen",
     "Debugging schwieriger;\nMessage Ordering"),
]

y_h = Inches(2.0)
txt(slide, Inches(0.8), y_h, Inches(2), Inches(0.3), "ENTSCHEIDUNG", size=10, bold=True, color=SUBTLE)
txt(slide, Inches(3.0), y_h, Inches(2), Inches(0.3), "LOEST NFR", size=10, bold=True, color=SUBTLE)
txt(slide, Inches(5.3), y_h, Inches(3.5), Inches(0.3), "VORTEILE", size=10, bold=True, color=SUBTLE)
txt(slide, Inches(9.5), y_h, Inches(3.5), Inches(0.3), "TRADE-OFFS", size=10, bold=True, color=SUBTLE)

y = Inches(2.5)
for decision, nfr, pro, con in decisions:
    txt(slide, Inches(0.8), y, Inches(2), Inches(1), decision, size=14, bold=True, color=ACCENT_TEAL)
    txt(slide, Inches(3.0), y, Inches(2.2), Inches(1), nfr, size=11, color=ACCENT_ORANGE)
    txt(slide, Inches(5.3), y, Inches(3.8), Inches(1), pro, size=11, color=LIGHT)
    txt(slide, Inches(9.5), y, Inches(3.5), Inches(1), con, size=11, color=RGBColor(0xCC, 0x88, 0x88))
    y += Inches(1.15)

add_notes(slide, """ENTSCHEIDUNGEN (4 Min):
- "Jede Entscheidung hat einen klaren Grund"
- Microservices: "Nicht weil es hip ist, sondern weil wir UNABHAENGIG deployen und skalieren muessen"
- Adapter: "DAS ist der Schluessel - wenn Keycloak kommt, tauschen wir nur den Adapter"
- Trade-offs EHRLICH ansprechen: "Microservices = Komplexitaet, aber beherrschbar"
- WICHTIG: Zeigt dass man Patterns UND deren Kosten kennt -> Senior-Denken!""")


# ═══════════════════════════════════════════════════════════
# SLIDE 6: HIGH-LEVEL ARCHITEKTUR
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, RGBColor(0xFA, 0xFA, 0xFC))
side_accent(slide)

txt(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
    "High-Level Architektur", size=32, bold=True, color=DARK_TEXT)
s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.0), Inches(2.5), Inches(0.04))
s.fill.solid()
s.fill.fore_color.rgb = ACCENT_RED
s.line.fill.background()

# Frontend
box(slide, Inches(3.5), Inches(1.3), Inches(6.5), Inches(0.65),
    "Angular SPA (Responsive / PWA)", ACCENT_RED, WHITE, 13,
    "Lazy Loading - Service Workers - SSR-ready")

# API Gateway
box(slide, Inches(3.5), Inches(2.2), Inches(6.5), Inches(0.6),
    "API Gateway  (Spring Cloud Gateway)", RGBColor(0x2C, 0x3E, 0x50), WHITE, 12,
    "Routing - Rate Limiting - Auth Validation - CORS")

# Services
svc_y = Inches(3.1)
svcs = [
    ("Article\nService", ACCENT_TEAL),
    ("User\nService", ACCENT_BLUE),
    ("Subscription\nService", PURPLE),
    ("Forum\nService", ACCENT_ORANGE),
    ("Archive\nService", RGBColor(0x5D, 0x6D, 0x7E)),
]
svc_w = Inches(2.2)
for i, (name, color) in enumerate(svcs):
    x = Inches(1.2) + i * (svc_w + Inches(0.25))
    box(slide, x, svc_y, svc_w, Inches(0.85), name, color, WHITE, 11)

# Data Layer
db_y = Inches(4.25)
box(slide, Inches(1.2), db_y, Inches(3.2), Inches(0.55),
    "PostgreSQL (DB per Service)", RGBColor(0x34, 0x49, 0x5E), WHITE, 10)
box(slide, Inches(4.6), db_y, Inches(2.2), Inches(0.55),
    "Redis (Cache / Sessions)", RGBColor(0x8E, 0x44, 0xAD), WHITE, 10)
box(slide, Inches(7.0), db_y, Inches(3.0), Inches(0.55),
    "Elasticsearch (Archiv + Suche)", RGBColor(0xF3, 0x9C, 0x12), WHITE, 10)
box(slide, Inches(10.2), db_y, Inches(2.2), Inches(0.55),
    "S3 (Media Assets)", RGBColor(0x1A, 0xBC, 0x9C), WHITE, 10)

# External systems
ext_y = Inches(5.4)
txt(slide, Inches(0.5), Inches(5.1), Inches(3), Inches(0.3),
    "EXTERNE SYSTEME (via Adapter)", size=10, bold=True, color=MID)

adapters_list = [
    "Keycloak\n(Auth/OIDC)", "Cloud CMS\n(Content)", "Buchhaltung\n(Legacy/SOAP)",
    "Merch-Shop\n(REST API)", "Payment\n(Stripe/PayPal)",
]
for i, name in enumerate(adapters_list):
    x = Inches(1.2) + i * Inches(2.45)
    a = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, ext_y - Inches(0.2), svc_w, Inches(0.18))
    a.fill.solid()
    a.fill.fore_color.rgb = ACCENT_ORANGE
    a.line.fill.background()
    box(slide, x, ext_y, svc_w, Inches(0.7), name, RGBColor(0xBD, 0xC3, 0xC7), DARK_TEXT, 9)

txt(slide, Inches(1.2), ext_y - Inches(0.4), Inches(4), Inches(0.2),
    "▼ Adapter Layer (Anti-Corruption)", size=9, bold=True, color=ACCENT_ORANGE)

txt(slide, Inches(0.5), Inches(6.8), Inches(12), Inches(0.4),
    "Orange = Adapter-Schicht (austauschbar)    |    Jeder Service eigene DB    |    Gateway = zentrale Sicherheit",
    size=10, color=MID)

add_notes(slide, """ARCHITEKTUR (5 Min) - KERNFOLIE:
Von oben nach unten:
1. "User sieht Angular-Frontend, responsive, PWA"
2. "Alle Requests durch API Gateway - zentrale Sicherheit"
3. "5 unabhaengige Microservices, jeder mit eigener DB"
4. "Und HIER der Clou:" -> orangene Adapter zeigen

"Die Adapter-Schicht ist das Herzstueck der Modifizierbarkeit:
 Wenn Keycloak statt AD kommt, tausche ich NUR den Auth-Adapter.
 Der Rest merkt NICHTS."

"Database per Service = kein shared state, lose Kopplung"
Bei Fragen zu Konsistenz: "Eventual Consistency ueber Events, Saga-Pattern wo noetig"

PAUSE: Fragen ob Rueckfragen da sind""")


# ═══════════════════════════════════════════════════════════
# SLIDE 7: ADAPTER PATTERN DETAIL
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Adapter Pattern im Detail", "Wie Austauschbarkeit technisch funktioniert")

txt(slide, Inches(1.0), Inches(2.0), Inches(5), Inches(0.4),
    "Beispiel: Auth-Adapter", size=18, bold=True, color=ACCENT_TEAL)

code_lines = [
    ("interface AuthPort {", {"size": 13, "color": ACCENT_ORANGE, "bold": True}),
    ("  authenticate(credentials): Token", {"size": 13, "color": LIGHT}),
    ("  validateToken(token): UserInfo", {"size": 13, "color": LIGHT}),
    ("  getUserRoles(userId): Role[]", {"size": 13, "color": LIGHT}),
    ("}", {"size": 13, "color": ACCENT_ORANGE, "bold": True}),
    ("", {"size": 8, "color": LIGHT}),
    ("// Heute:", {"size": 13, "color": MID}),
    ("class ADAuthAdapter implements AuthPort", {"size": 13, "color": GREEN, "bold": True}),
    ("", {"size": 8, "color": LIGHT}),
    ("// Morgen (wenn entschieden):", {"size": 13, "color": MID}),
    ("class KeycloakAdapter implements AuthPort", {"size": 13, "color": ACCENT_BLUE, "bold": True}),
]
multi_txt(slide, Inches(1.0), Inches(2.5), Inches(5.5), Inches(4.5), code_lines, size=13, color=LIGHT, spacing=Pt(2))

txt(slide, Inches(7.5), Inches(2.0), Inches(5), Inches(0.4),
    "Warum das funktioniert:", size=18, bold=True, color=ACCENT_ORANGE)

benefits = [
    ("✓  Dependency Inversion", {"size": 15, "color": WHITE, "bold": True, "spacing": Pt(4)}),
    ("    Services kennen nur das Interface", {"size": 12, "color": SUBTLE, "spacing": Pt(14)}),
    ("✓  Konfigurationsgetrieben", {"size": 15, "color": WHITE, "bold": True, "spacing": Pt(4)}),
    ("    Adapter-Wechsel via Config/Feature Flag", {"size": 12, "color": SUBTLE, "spacing": Pt(14)}),
    ("✓  Testbar", {"size": 15, "color": WHITE, "bold": True, "spacing": Pt(4)}),
    ("    Mock-Adapter fuer Tests, kein Ext.-System noetig", {"size": 12, "color": SUBTLE, "spacing": Pt(14)}),
    ("✓  Parallel betreibbar", {"size": 15, "color": WHITE, "bold": True, "spacing": Pt(4)}),
    ("    AD + Keycloak gleichzeitig (Canary Migration)", {"size": 12, "color": SUBTLE, "spacing": Pt(14)}),
]
multi_txt(slide, Inches(7.5), Inches(2.5), Inches(5), Inches(4.5), benefits, size=15, color=WHITE, spacing=Pt(4))

add_notes(slide, """ADAPTER DETAIL (2 Min):
- "Konkret, wie das Pattern funktioniert"
- Links: "Interface definiert WAS wir brauchen - nicht WIE"
- "Heute AD-Adapter. Keycloak kommt? Neuer Adapter, selbes Interface"
- BETONEN: "Beide PARALLEL betreibbar - Canary-Migration, kein Big-Bang"
- Zeigt: Migrationspfad durchdacht, nicht nur Zielarchitektur""")


# ═══════════════════════════════════════════════════════════
# SLIDE 8: TECH STACK
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Technology Stack", "Bewaehrte Technologien, Cloud-native Ausrichtung")

categories = [
    ("Frontend", ["Angular 18+", "TypeScript", "SCSS + Responsive", "PWA-faehig"], ACCENT_RED),
    ("Backend", ["Java 21 + Spring Boot 3", "Spring Cloud Gateway", "Spring Security", "OpenAPI 3.0"], ACCENT_TEAL),
    ("Data", ["PostgreSQL (per Service)", "Redis (Cache)", "Elasticsearch", "S3 (Media)"], ACCENT_BLUE),
    ("Infra", ["AWS ECS Fargate", "Docker", "Terraform (IaC)", "GitHub Actions"], PURPLE),
    ("Monitoring", ["Prometheus + Grafana", "ELK Stack", "Distributed Tracing", "Alerting"], ACCENT_ORANGE),
]

x = Inches(0.6)
for cat_name, items, color in categories:
    txt(slide, x, Inches(2.0), Inches(2.3), Inches(0.35), cat_name, size=13, bold=True, color=color)
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(2.35), Inches(2.0), Inches(0.03))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    y = Inches(2.55)
    for item in items:
        txt(slide, x, y, Inches(2.4), Inches(0.3), f"▸ {item}", size=11, color=LIGHT)
        y += Inches(0.38)
    x += Inches(2.5)

add_notes(slide, """TECH STACK (1-2 Min):
- Kurz durchgehen
- "Angular + Spring Boot - bewaehrt, grosses Oekosystem, gute Entwickler-Verfuegbarkeit"
- "Jeder Service eigene PostgreSQL-DB"
- "AWS Fargate - serverless Container, kein Cluster-Management"
- "Observability von Tag 1 - bei Microservices MUSS man wissen was passiert" """)


# ═══════════════════════════════════════════════════════════
# SLIDE 9: EVOLUTIONSPFAD
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Evolutionspfad", "Iterativer Ausbau - kein Big-Bang")

phases = [
    ("Phase 1\nMVP", "Q3 2026", ["Article + User Service", "Basic Auth (AD-Adapter)", "Public-Bereich live", "CI/CD Pipeline"], ACCENT_TEAL),
    ("Phase 2\nMonetarisierung", "Q4 2026", ["Subscription Service", "Payment-Integration", "Premium-Bereich", "Abo-Management"], ACCENT_BLUE),
    ("Phase 3\nCommunity", "Q1 2027", ["Forum Service", "Moderation-Tools", "Benachrichtigungen", "Author-Dashboard"], PURPLE),
    ("Phase 4\nOptimierung", "Q2 2027", ["Archive + Elasticsearch", "Keycloak-Migration", "CMS-Abloesung", "Performance-Tuning"], ACCENT_ORANGE),
]

s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(3.0), Inches(12.0), Inches(0.03))
s.fill.solid()
s.fill.fore_color.rgb = MID
s.line.fill.background()

x = Inches(0.6)
for phase_name, timeline, items, color in phases:
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(1.2), Inches(2.85), Inches(0.3), Inches(0.3))
    c.fill.solid()
    c.fill.fore_color.rgb = color
    c.line.fill.background()
    txt(slide, x, Inches(2.0), Inches(2.8), Inches(0.7), phase_name, size=15, bold=True, color=color)
    txt(slide, x, Inches(3.3), Inches(2.5), Inches(0.3), timeline, size=11, color=SUBTLE)
    y = Inches(3.7)
    for item in items:
        txt(slide, x, y, Inches(2.8), Inches(0.3), f"· {item}", size=11, color=LIGHT)
        y += Inches(0.35)
    x += Inches(3.1)

box(slide, Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.7),
    "Jede Phase liefert Mehrwert · Parallele Teamarbeit moeglich · Migrationen Feature-Flag-gesteuert",
    RGBColor(0x1A, 0x2A, 0x1A), GREEN, size=12)

add_notes(slide, """EVOLUTIONSPFAD (3 Min):
- Phase 1: "In 2-3 Monaten lesbares Magazin live - echter User-Value"
- Phase 2: "Monetarisierung - Kunde verdient Geld"
- Phase 3: "Community fuer Retention"
- Phase 4: "HIER die unsicheren Migrationen - bis dahin wissen wir mehr"
- BETONEN: "Adapter-Architektur = FREIHEIT, Migrationen zu verschieben ohne tech debt"
- "Agile Architektur: Nicht alles planen, aber auf alles VORBEREITET sein" """)


# ═══════════════════════════════════════════════════════════
# SLIDE 10: RISIKEN
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Risiken & Gegenmassnahmen", "Proaktives Risikomanagement")

risks = [
    ("HOCH", "Legacy-Buchhaltung (Win2000)\nkeine moderne API", "SOAP/REST-Wrapper;\nAsync Message Queue als Puffer", ACCENT_RED),
    ("HOCH", "CMS-Entscheidung offen\nTeam plant ohne Klarheit", "Content-Adapter von Tag 1;\nFeature Flags fuer Switch", ACCENT_RED),
    ("MITTEL", "Microservice-Komplexitaet\nfuer kleines Team", "Start mit 2-3 Services;\nGraduelle Extraktion moeglich", ACCENT_ORANGE),
    ("MITTEL", "Lastspitzen unvorhersehbar\n(virale Events)", "Auto-Scaling; CDN;\nCircuit Breaker Pattern", ACCENT_ORANGE),
]

y = Inches(2.0)
txt(slide, Inches(0.8), y, Inches(1), Inches(0.3), "RISIKO", size=10, bold=True, color=SUBTLE)
txt(slide, Inches(2.0), y, Inches(3.5), Inches(0.3), "BESCHREIBUNG", size=10, bold=True, color=SUBTLE)
txt(slide, Inches(6.5), y, Inches(5.5), Inches(0.3), "MITIGATION", size=10, bold=True, color=SUBTLE)

y = Inches(2.4)
for level, desc, mitigation, color in risks:
    box(slide, Inches(0.8), y + Inches(0.05), Inches(0.9), Inches(0.35), level, color, WHITE, 9)
    txt(slide, Inches(2.0), y, Inches(4.2), Inches(0.9), desc, size=12, color=LIGHT)
    txt(slide, Inches(6.5), y, Inches(6.0), Inches(0.9), mitigation, size=12, color=GREEN)
    y += Inches(1.05)
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), y - Inches(0.1), Inches(11.5), Inches(0.01))
    s.fill.solid()
    s.fill.fore_color.rgb = RGBColor(0x2A, 0x2A, 0x3E)
    s.line.fill.background()

add_notes(slide, """RISIKEN (2 Min):
- "Kein Projekt ohne Risiken - die 4 wichtigsten:"
- Legacy: "Win2000 wird nicht ersetzt - Wrapper drumherum"
- CMS: "Genau deshalb der Adapter"
- Komplexitaet: "Wir muessen nicht mit 5 Services STARTEN. Modular Monolith Start moeglich!"
  -> DAS ZEIGT PRAGMATISMUS!
- Lastspitzen: "Auto-Scaling + CDN"
- Zeigt: Architektur UND Projektrisiken im Blick""")


# ═══════════════════════════════════════════════════════════
# SLIDE 11: KEY TAKEAWAYS
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)
side_accent(slide)

txt(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8), "Key Takeaways", size=36, bold=True, color=WHITE)
accent_bar(slide)

takeaways = [
    ("1", "Modifizierbarkeit als primaerer Architektur-Treiber", "Adapter Pattern + Dependency Inversion fuer alle Umsysteme"),
    ("2", "Microservices fuer unabhaengige Skalierung & Deployment", "Database per Service, API Gateway, Resilience Patterns"),
    ("3", "Iterativer Aufbau statt Big-Bang", "MVP in Q3, Monetarisierung in Q4 - jede Phase liefert Value"),
    ("4", "Risiken aktiv gemanagt", "Legacy-Wrapper, Feature Flags, gradueller Migrationspfad"),
    ("5", "Cloud-native & Observable von Tag 1", "AWS Fargate, Terraform IaC, Prometheus/Grafana Monitoring"),
]

y = Inches(1.9)
for num, title, desc in takeaways:
    txt(slide, Inches(0.8), y, Inches(0.6), Inches(0.5), num, size=24, bold=True, color=ACCENT_RED, font="Segoe UI Light")
    txt(slide, Inches(1.5), y, Inches(10), Inches(0.4), title, size=17, bold=True, color=WHITE)
    txt(slide, Inches(1.5), y + Inches(0.37), Inches(10), Inches(0.3), desc, size=12, color=SUBTLE)
    y += Inches(0.9)

txt(slide, Inches(0.8), Inches(6.5), Inches(10), Inches(0.5),
    "Vielen Dank! Bereit fuer Fragen & Diskussion", size=20, bold=True, color=ACCENT_TEAL)

add_notes(slide, """ABSCHLUSS (1 Min):
- 5 Takeaways je 1 Satz
- Selbstbewusst: "Die Architektur ist auf Veraenderung ausgelegt - genau wie ein agiles Projekt es braucht."
- "Ich freue mich auf Ihre Fragen!"

BEI TYPISCHEN FRAGEN:
- "Warum nicht Monolith?" -> "Start KANN modular-monolithisch sein, Extraktion spaeter"
- "Kosten?" -> "Fargate = pay per use, skaliert mit Bedarf"
- "Team-Groesse?" -> "2-3 Devs starten mit 2 Services, waechst mit Projekt"
- "Warum nicht K8s?" -> "Fargate = weniger Ops-Overhead, besser fuer kleines Team"
- "Wie testet man?" -> "Contract Tests zwischen Services, Mock-Adapter fuer externe Systeme"
- "Event-Driven vs REST?" -> "REST fuer synchrone Reads, Events fuer async Writes"

GOLDEN RULE: Bei jeder Antwort ABWAEGEN zeigen, nicht dogmatisch sein!""")


# ═══════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════
output = "/Users/gonibeer/dev/extremsport-magazine/docs/Case_Study_Extremsport_Magazin.pptx"
prs.save(output)
print(f"✅ Presentation saved: {output}")
print(f"   📊 {len(prs.slides)} slides")
print(f"   🎤 Speaker notes on every slide (timing + delivery tips)")
print(f"   ⏱  Designed for ~20 minutes")
print(f"   💡 Q&A preparation included in final slide notes")

